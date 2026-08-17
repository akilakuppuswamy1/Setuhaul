"""Build SetuHaul FDE assignment PowerPoint (widescreen 16:9)."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

DOCS = Path(__file__).resolve().parent
OUT = DOCS / "SetuHaul_FDE_Assignment_Presentation.pptx"

NAVY = RGBColor(0x0F, 0x27, 0x44)
NAVY_2 = RGBColor(0x12, 0x30, 0x56)
GOLD = RGBColor(0xC4, 0x8A, 0x1A)
GOLD_LT = RGBColor(0xF4, 0xE4, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xEE, 0xF2, 0xF6)
MUTED = RGBColor(0x4A, 0x5D, 0x73)
INK = RGBColor(0x10, 0x24, 0x3E)
GREEN = RGBColor(0x1F, 0x6B, 0x4A)
GREEN_BG = RGBColor(0xE4, 0xF3, 0xEA)
RED = RGBColor(0x8B, 0x2E, 0x4A)
RED_BG = RGBColor(0xF7, 0xE8, 0xEF)
BLUE = RGBColor(0x2F, 0x5F, 0x8A)
BLUE_BG = RGBColor(0xE3, 0xED, 0xF6)
AMBER_BG = RGBColor(0xFF, 0xF6, 0xE8)
SLATE = RGBColor(0x5C, 0x6B, 0x7A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _set_run(run, text, size=14, bold=False, color=INK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _fill_line(shape, fill, line=None, width_pt=1.0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(width_pt)


def box(
    slide,
    l,
    t,
    w,
    h,
    fill,
    text="",
    size=12,
    bold=False,
    color=WHITE,
    align=PP_ALIGN.CENTER,
    rounded=True,
    line=None,
    anchor=MSO_ANCHOR.MIDDLE,
):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape, l, t, w, h)
    _fill_line(s, fill, line)
    if rounded:
        s.adjustments[0] = 0.08
    tf = s.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        s.text_frame.paragraphs[0].alignment = align
    except Exception:
        pass
    tf.paragraphs[0].alignment = align
    s.text_frame.paragraphs[0].space_before = Pt(0)
    s.text_frame.paragraphs[0].space_after = Pt(0)
    # vertical anchor
    body = s._element.find(qn("p:txBody"))
    if body is not None:
        bodyPr = body.find(qn("a:bodyPr"))
        if bodyPr is not None:
            bodyPr.set("anchor", {MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.TOP: "t", MSO_ANCHOR.BOTTOM: "b"}[anchor])
            bodyPr.set("lIns", str(int(Emu(Inches(0.08)))))
            bodyPr.set("rIns", str(int(Emu(Inches(0.08)))))
            bodyPr.set("tIns", str(int(Emu(Inches(0.04)))))
            bodyPr.set("bIns", str(int(Emu(Inches(0.04)))))
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size=size, bold=bold, color=color)
    return s


def textbox(slide, l, t, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    _set_run(run, text, size=size, bold=bold, color=color)
    return tb


def add_lines(slide, l, t, w, h, lines, size=14, color=INK, bold_first=False, spacing=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        run = p.add_run()
        _set_run(run, line, size=size, bold=(bold_first and i == 0), color=color)
    return tb


def footer(slide, n, total):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.28), SLIDE_W, Inches(0.22))
    _fill(bar, NAVY)
    textbox(slide, Inches(0.35), Inches(7.28), Inches(8), Inches(0.22), "SETUHAUL  ·  FDE Assignment", 10, False, GOLD_LT)
    textbox(
        slide,
        Inches(11.4),
        Inches(7.28),
        Inches(1.6),
        Inches(0.22),
        f"{n}  /  {total}",
        10,
        False,
        WHITE,
        PP_ALIGN.RIGHT,
    )


def header(slide, kicker, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(0.92))
    _fill(band, NAVY)
    gold = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.92), SLIDE_W, Inches(0.06))
    _fill(gold, GOLD)
    textbox(slide, Inches(0.4), Inches(0.08), Inches(12.4), Inches(0.28), kicker.upper(), 11, True, GOLD)
    textbox(slide, Inches(0.4), Inches(0.34), Inches(12.4), Inches(0.5), title, 24, True, WHITE)


def arrow_right(slide, l, t, w=Inches(0.32), h=Inches(0.22)):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, w, h)
    _fill(s, GOLD)
    return s


def arrow_down(slide, l, t, w=Inches(0.22), h=Inches(0.28)):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, l, t, w, h)
    _fill(s, GOLD)
    return s


def _crop_png(name: str, y0: float = 0.0, y1: float = 1.0) -> tuple[BytesIO, tuple[int, int]]:
    im = Image.open(DOCS / name).convert("RGB")
    w, h = im.size
    crop = im.crop((0, int(h * y0), w, max(int(h * y0) + 1, int(h * y1))))
    buf = BytesIO()
    crop.save(buf, format="PNG")
    buf.seek(0)
    return buf, crop.size


def place_fitted(slide, name: str, left, top, max_w, max_h, y0: float = 0.0, y1: float = 1.0):
    buf, (iw, ih) = _crop_png(name, y0, y1)
    scale = min(float(max_w) / iw, float(max_h) / ih)
    pw, ph = int(iw * scale), int(ih * scale)
    l = int(left + (max_w - pw) / 2)
    t = int(top + (max_h - ph) / 2)
    slide.shapes.add_picture(buf, l, t, pw, ph)


def diagram_slide(add, kicker: str, title: str, png: str, y0: float = 0.0, y1: float = 1.0, caption: str | None = None):
    s = add()
    header(s, kicker, title)
    frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(1.1), Inches(12.77), Inches(6.05 if caption else 6.12))
    _fill_line(frame, WHITE, RGBColor(0xD5, 0xDE, 0xE8), 1.0)
    img_top = Inches(1.16)
    img_h = Inches(5.55 if caption else 5.98)
    place_fitted(s, png, Inches(0.36), img_top, Inches(12.6), img_h, y0, y1)
    if caption:
        textbox(s, Inches(0.4), Inches(6.78), Inches(12.5), Inches(0.38), caption, 12, False, MUTED)
    return s


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    prs = new_prs()
    slides = []

    def add():
        s = blank(prs)
        slides.append(s)
        return s

    # ---------- 1 Title ----------
    s = add()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _fill(bg, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    _fill(accent, GOLD)
    textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.35), "FDE CHALLENGE  ·  ASSIGNMENT PRESENTATION", 14, True, GOLD)
    textbox(s, Inches(0.7), Inches(1.85), Inches(12), Inches(1.2), "SETUHAUL", 54, True, WHITE)
    textbox(
        s,
        Inches(0.7),
        Inches(3.05),
        Inches(11.5),
        Inches(0.9),
        "Deterministic warehouse appointment coordination\nwith conversational AI",
        22,
        False,
        GOLD_LT,
    )
    box(s, Inches(0.7), Inches(4.35), Inches(3.4), Inches(0.85), RGBColor(0x1A, 0x3A, 0x5C), "LLM for language\nPython for decisions", 13, True, WHITE)
    box(s, Inches(4.3), Inches(4.35), Inches(3.4), Inches(0.85), RGBColor(0x1A, 0x3A, 0x5C), "Show → Propose → Confirm\ncapacity only at commit", 13, True, WHITE)
    box(s, Inches(7.9), Inches(4.35), Inches(4.4), Inches(0.85), RGBColor(0x1A, 0x3A, 0x5C), "Steps 1–9 implemented\n399 backend tests", 13, True, WHITE)
    textbox(s, Inches(0.7), Inches(5.55), Inches(11), Inches(0.7), "Classroom problem: late inbound driver + scarce concurrent capacity\nOptional: read-only facility ranking (Step 9) — never books", 14, False, RGBColor(0xC5, 0xD0, 0xDC))
    textbox(s, Inches(0.7), Inches(6.55), Inches(11), Inches(0.35), "No LangChain  ·  No LangGraph  ·  No OR-Tools  ·  Frozen 16-table schema", 13, False, GOLD)

    # ---------- 2 Agenda ----------
    s = add()
    header(s, "Overview", "Agenda")
    items = [
        ("01", "Problem", "Scarce docks, late drivers, unsafe naive booking"),
        ("02", "Architecture", "AI boundary, layered stack, frozen schema"),
        ("03", "Authority path", "Feasibility → proposal → locked allocation"),
        ("04", "Conversation", "Allowlisted tools over SH-1024 hero flow"),
        ("05", "Safety", "Concurrency, stale proposals, hardening"),
        ("06", "Evidence", "Tests, demo console, limitations"),
    ]
    for i, (num, title, blurb) in enumerate(items):
        col, row = i % 3, i // 3
        l = Inches(0.45 + col * 4.2)
        t = Inches(1.35 + row * 2.7)
        box(s, l, t, Inches(3.95), Inches(2.4), SOFT, "", color=INK, line=RGBColor(0xD5, 0xDE, 0xE8))
        box(s, l, t, Inches(3.95), Inches(0.55), NAVY, num + "   " + title, 16, True, WHITE, rounded=False)
        add_lines(s, l + Inches(0.18), t + Inches(0.75), Inches(3.55), Inches(1.4), [blurb], size=15, color=INK)

    # ---------- 3 Problem ----------
    s = add()
    header(s, "Context", "The receiving dock cannot treat chat as a booking")
    pains = [
        ("Late arrival", "Original 6:30 PM window.\nNew ETA 8:30 PM."),
        ("Leave-by constraint", "Driver must leave by 9:30 PM.\nLater slots are useless."),
        ("Scarce capacity", "Finite docks and slot counts.\nSeveral trucks want the last window."),
        ("Stale options", "What was feasible when shown\nmay be gone at confirm time."),
        ("Race to confirm", "Two clients both see remaining\ncapacity, then both confirm."),
        ("Language ≠ authority", "A chat message is not a lock,\nnot a hold, not a booking."),
    ]
    for i, (title, body) in enumerate(pains):
        col, row = i % 3, i // 3
        l = Inches(0.4 + col * 4.25)
        t = Inches(1.25 + row * 2.85)
        box(s, l, t, Inches(4.05), Inches(2.6), WHITE, "", line=RGBColor(0xD0, 0xDA, 0xE4))
        box(s, l, t, Inches(0.12), Inches(2.6), GOLD if row == 0 else RED, rounded=False)
        textbox(s, l + Inches(0.3), t + Inches(0.25), Inches(3.55), Inches(0.5), title, 18, True, NAVY)
        add_lines(s, l + Inches(0.3), t + Inches(0.85), Inches(3.55), Inches(1.5), body.split("\n"), size=15, color=MUTED)

    # ---------- 4 Naive booking ----------
    s = add()
    header(s, "Failure mode", "Naive “book the slot I asked for” is unsafe")
    add_lines(
        s,
        Inches(0.45),
        Inches(1.15),
        Inches(12.4),
        Inches(0.4),
        ["Showing an option is not holding capacity. Feasibility at display time is not feasibility at commit time."],
        size=16,
        color=MUTED,
    )
    # left flow
    box(s, Inches(0.4), Inches(1.7), Inches(6.2), Inches(5.2), RED_BG, rounded=False)
    textbox(s, Inches(0.6), Inches(1.85), Inches(5.8), Inches(0.4), "UNSAFE PATH", 14, True, RED)
    steps_bad = [
        "Driver: “Book me 8:30 PM”",
        "Chat treats message as booking",
        "No re-check at commit",
        "No row locks on slot / dock",
        "Two confirms both succeed",
        "Double-booked last unit of capacity",
    ]
    for i, t in enumerate(steps_bad):
        box(s, Inches(0.7), Inches(2.35 + i * 0.7), Inches(5.6), Inches(0.55), WHITE, t, 13, False, INK, line=RED)
        if i < len(steps_bad) - 1:
            arrow_down(s, Inches(3.3), Inches(2.88 + i * 0.7), Inches(0.18), Inches(0.16))

    box(s, Inches(6.85), Inches(1.7), Inches(6.05), Inches(5.2), GREEN_BG, rounded=False)
    textbox(s, Inches(7.05), Inches(1.85), Inches(5.65), Inches(0.4), "SETUHAUL PATH", 14, True, GREEN)
    steps_good = [
        "Record facts (ETA, leave-by)",
        "Evaluate with Step 5 rules",
        "Propose without consuming capacity",
        "Explicit confirm language only",
        "Revalidate + lock + capacity re-check",
        "One winner, loser gets HTTP 409 / stale",
    ]
    for i, t in enumerate(steps_good):
        box(s, Inches(7.15), Inches(2.35 + i * 0.7), Inches(5.55), Inches(0.55), WHITE, t, 13, False, INK, line=GREEN)

    # ---------- 5 Assignment mapping ----------
    s = add()
    header(s, "Assignment", "Classroom problem mapped onto this repository")
    rows = [
        ("Driver conversation", "Step 8 over frozen ChatThread / ChatMessage"),
        ("Scarce appointment capacity", "Slot capacity, docks, facility rules"),
        ("Feasibility", "Step 5 FeasibilityEngine — pure rules, no SQL"),
        ("Controlled proposals", "Step 7 Appointment rows with status=requested"),
        ("Concurrency", "Step 6 AllocationService: locks + capacity re-check"),
        ("Optional facility scheduling", "Step 9 read-only ranking — no confirm endpoint"),
    ]
    box(s, Inches(0.4), Inches(1.2), Inches(4.4), Inches(0.5), NAVY, "Theme", 14, True, WHITE, rounded=False)
    box(s, Inches(4.8), Inches(1.2), Inches(8.1), Inches(0.5), GOLD, "How SetuHaul treats it", 14, True, NAVY, rounded=False)
    for i, (a, b) in enumerate(rows):
        bg = SOFT if i % 2 == 0 else WHITE
        top = Inches(1.72 + i * 0.78)
        box(s, Inches(0.4), top, Inches(4.4), Inches(0.72), bg, a, 14, True, NAVY, rounded=False, line=RGBColor(0xD8, 0xE0, 0xE8))
        box(s, Inches(4.8), top, Inches(8.1), Inches(0.72), bg, b, 14, False, INK, rounded=False, line=RGBColor(0xD8, 0xE0, 0xE8))
    textbox(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.45), "Required path: Steps 5–8.  Step 9 never books, locks, reserves, or confirms capacity.", 13, True, MUTED)

    # ---------- 6 Principle ----------
    s = add()
    header(s, "Core design", "LLM for language. Deterministic Python for decisions.")
    box(s, Inches(0.4), Inches(1.25), Inches(6.1), Inches(5.7), GREEN_BG, rounded=False)
    textbox(s, Inches(0.6), Inches(1.4), Inches(5.7), Inches(0.4), "WHAT AI DOES", 16, True, GREEN)
    does = [
        "Understand driver language (intent, times, option ordinals)",
        "Classify intent and extract ETA / leave-by / earliest start",
        "Keep conversation context in ChatMessage.metadata JSON",
        "Choose allowlisted tools only",
        "Explain deterministic rule_results in driver-facing text",
        "Request human escalation as a record / flag",
    ]
    for i, t in enumerate(does):
        box(s, Inches(0.65), Inches(1.95 + i * 0.75), Inches(5.6), Inches(0.65), WHITE, "  " + t, 13, False, INK, PP_ALIGN.LEFT)

    box(s, Inches(6.75), Inches(1.25), Inches(6.15), Inches(5.7), RED_BG, rounded=False)
    textbox(s, Inches(6.95), Inches(1.4), Inches(5.75), Inches(0.4), "WHAT AI DOES NOT DO", 16, True, RED)
    dont = [
        "Decide feasibility or invent travel times",
        "Act as capacity authority",
        "Choose / lock docks as a commit",
        "Allocate or confirm independently",
        "Bypass concurrency controls or write SQL",
        "Override Step 5 / 6 / 7 rules",
    ]
    for i, t in enumerate(dont):
        box(s, Inches(7.0), Inches(1.95 + i * 0.75), Inches(5.65), Inches(0.65), WHITE, "  " + t, 13, False, INK, PP_ALIGN.LEFT)

    # ---------- 7 Required path ----------
    s = add()
    header(s, "Solution", "Required confirmation path (driver → booking)")
    nodes = [
        ("Driver", NAVY),
        ("React\nconsole", BLUE),
        ("FastAPI", BLUE),
        ("Step 8\nconversation", RGBColor(0x3D, 0x5A, 0x80)),
        ("Step 5\nfeasibility", GREEN),
        ("Step 7\nproposal", GOLD),
        ("Step 5\nrevalidate", GREEN),
        ("Step 6\nallocate", RED),
        ("Confirmed\nappointment", NAVY),
    ]
    y = Inches(2.15)
    for i, (label, color) in enumerate(nodes):
        x = Inches(0.28 + i * 1.45)
        box(s, x, y, Inches(1.32), Inches(1.15), color, label, 11, True, WHITE)
        if i < len(nodes) - 1:
            arrow_right(s, Inches(1.55 + i * 1.45), Inches(2.55), Inches(0.22), Inches(0.28))
    box(
        s,
        Inches(0.4),
        Inches(3.7),
        Inches(12.5),
        Inches(0.7),
        AMBER_BG,
        "Privileged confirmation is only:  Step 8 accept_proposal  →  Step 7 accept  →  Step 5 revalidation  →  Step 6 allocation",
        14,
        True,
        NAVY,
        line=GOLD,
    )
    # optional path
    textbox(s, Inches(0.45), Inches(4.6), Inches(12), Inches(0.35), "Optional path — Step 9 (no booking)", 16, True, NAVY)
    opt = [
        ("Facility\nsnapshot", BLUE),
        ("Scheduling\nService", BLUE),
        ("Step 5\neligibility", GREEN),
        ("Scheduling\nEngine rank", GOLD),
        ("Proposed\nschedule", NAVY),
        ("NO booking\nNO lock", RED),
    ]
    for i, (label, color) in enumerate(opt):
        x = Inches(0.4 + i * 2.15)
        box(s, x, Inches(5.1), Inches(1.95), Inches(1.05), color, label, 12, True, WHITE)
        if i < len(opt) - 1:
            arrow_right(s, Inches(2.32 + i * 2.15), Inches(5.48), Inches(0.22), Inches(0.26))

    # ---------- 8 Layered architecture ----------
    s = add()
    header(s, "Architecture", "Layered stack — UI never calculates feasibility")
    layers = [
        ("Clients", "React 19 + Vite ops console  ·  Driver, Shipments, Appointments, Facility Schedule, Demo, Concurrency", RGBColor(0x1A, 0x3A, 0x5C)),
        ("HTTP API", "FastAPI  ·  Pydantic contracts  ·  CORS  ·  GET /health {service: setuhaul}", BLUE),
        ("Conversation", "Intent, entities, thread context, allowlisted tools, formatter  ·  FakeLLM default", RGBColor(0x3D, 0x5A, 0x80)),
        ("Orchestration", "Conversation / Feasibility / Proposal / Allocation / Scheduling services — load facts, never invent them", GOLD),
        ("Decision cores", "FeasibilityEngine (pure)  ·  Allocation locks (transactional)  ·  SchedulingEngine (rank only)", GREEN),
        ("Persistence", "PostgreSQL 16  ·  SQLAlchemy 2  ·  Alembic  ·  Frozen 16-table Step 2 schema", NAVY),
    ]
    for i, (name, desc, color) in enumerate(layers):
        top = Inches(1.2 + i * 0.92)
        box(s, Inches(0.4), top, Inches(2.4), Inches(0.8), color, name, 14, True, WHITE, rounded=False)
        box(s, Inches(2.8), top, Inches(10.1), Inches(0.8), SOFT, "  " + desc, 14, False, INK, PP_ALIGN.LEFT, rounded=False, line=color)
        if i < len(layers) - 1:
            arrow_down(s, Inches(1.4), Inches(1.98 + i * 0.92), Inches(0.2), Inches(0.16))

    diagram_slide(
        add, "Diagram  ·  docs/system_architecture.png",
        "Implemented layered stack",
        "system_architecture.png", 0.00, 0.26,
        caption="Clients → FastAPI → services → engines → frozen PostgreSQL. UI has no decision logic.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/system_architecture.png",
        "Runtime request path and confirmation sequence",
        "system_architecture.png", 0.24, 0.58,
        caption="Free text and REST share Step 5 / 6–7. Confirmation revalidates, then allocates under locks.",
    )

    # ---------- 9 AI boundary diagram ----------
    s = add()
    header(s, "AI boundary", "Allowlisted tools are the only crossing")
    box(s, Inches(0.35), Inches(1.2), Inches(4.0), Inches(5.7), GREEN_BG, rounded=False)
    textbox(s, Inches(0.5), Inches(1.35), Inches(3.7), Inches(0.4), "LANGUAGE SIDE", 14, True, GREEN)
    lang = ["Parse free text", "Ask when facts missing", "Explain rule_results", "Flag human escalation"]
    for i, t in enumerate(lang):
        box(s, Inches(0.55), Inches(1.9 + i * 1.15), Inches(3.6), Inches(0.95), WHITE, t, 15, True, NAVY, line=GREEN)

    box(s, Inches(4.55), Inches(1.2), Inches(4.15), Inches(5.7), AMBER_BG, rounded=False)
    textbox(s, Inches(4.7), Inches(1.35), Inches(3.85), Inches(0.4), "BOUNDARY", 14, True, GOLD)
    box(s, Inches(4.75), Inches(2.0), Inches(3.75), Inches(1.3), GOLD, "Allowlisted tools only\nunknown name → forbidden", 14, True, NAVY)
    box(s, Inches(4.75), Inches(3.5), Inches(3.75), Inches(1.15), WHITE, "ToolArguments extra=forbid\nUUID validation", 13, True, INK, line=GOLD)
    box(s, Inches(4.75), Inches(4.85), Inches(3.75), Inches(1.15), WHITE, "Injection markers skip\nirreversible tools", 13, True, INK, line=GOLD)

    box(s, Inches(8.9), Inches(1.2), Inches(4.05), Inches(5.7), RED_BG, rounded=False)
    textbox(s, Inches(9.05), Inches(1.35), Inches(3.75), Inches(0.4), "AUTHORITY SIDE", 14, True, RED)
    auth = [
        "Step 5 FeasibilityEngine",
        "Step 7 ProposalService",
        "Step 6 AllocationService",
        "Step 9 SchedulingEngine",
        "ETA / exception facts",
    ]
    for i, t in enumerate(auth):
        box(s, Inches(9.1), Inches(1.85 + i * 0.95), Inches(3.65), Inches(0.82), WHITE, t, 13, True, NAVY, line=RED)

    diagram_slide(
        add, "Diagram  ·  docs/ai_responsibility_boundary.png",
        "Language side vs authority side",
        "ai_responsibility_boundary.png", 0.00, 0.44,
        caption="Allowlisted tools only. FakeLLM / OpenRouter parse; engines decide. No LangChain / LangGraph.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/ai_responsibility_boundary.png",
        "No AI booking pathway",
        "ai_responsibility_boundary.png", 0.40, 0.78,
        caption="Show → propose → confirm. Step 9 ranking is read-only and never a hold.",
    )

    # ---------- 10 Providers ----------
    s = add()
    header(s, "Language adapters", "Two providers. Neither can authorize confirmation.")
    box(s, Inches(0.4), Inches(1.25), Inches(6.15), Inches(4.4), SOFT, rounded=False)
    box(s, Inches(0.4), Inches(1.25), Inches(6.15), Inches(0.6), NAVY, "FakeLLMProvider  —  default", 16, True, WHITE, rounded=False)
    add_lines(
        s,
        Inches(0.65),
        Inches(2.05),
        Inches(5.7),
        Inches(3.3),
        [
            "Deterministic parser in intents.py",
            "No network calls",
            "Safe classroom default (LLM_PROVIDER=fake)",
            "Same tool plan as the live hero flow",
            "Confirm / reject flags come from the parser",
        ],
        size=16,
        color=INK,
        spacing=10,
    )
    box(s, Inches(6.75), Inches(1.25), Inches(6.15), Inches(4.4), SOFT, rounded=False)
    box(s, Inches(6.75), Inches(1.25), Inches(6.15), Inches(0.6), BLUE, "OpenRouterProvider  —  optional", 16, True, WHITE, rounded=False)
    add_lines(
        s,
        Inches(7.0),
        Inches(2.05),
        Inches(5.7),
        Inches(3.3),
        [
            "Language adapter only",
            "Cannot independently authorize confirmation",
            "Write intents the parser did not also see are rejected",
            "confirm / reject come from the parser, not the model payload",
            "Missing API key falls back to FakeLLM",
        ],
        size=16,
        color=INK,
        spacing=10,
    )
    box(
        s,
        Inches(0.4),
        Inches(5.85),
        Inches(12.5),
        Inches(1.05),
        AMBER_BG,
        "LangChain and LangGraph are not used. There is no AI booking pathway.\nOpenRouter cannot promote a non-write parser intent to ACCEPT_PROPOSAL.",
        15,
        True,
        NAVY,
        line=GOLD,
    )

    # ---------- 11 Schema ----------
    s = add()
    header(s, "System of record", "Frozen Step 2 schema — 16 tables, zero added later")
    clusters = [
        ("Actors", ["carriers", "drivers", "vehicles", "contacts"], BLUE),
        ("Move", ["shipments", "eta_updates", "driver_exceptions"], GREEN),
        ("Facility", ["facilities", "docks", "facility_rules", "appointment_slots"], GOLD),
        ("Commitment", ["appointments", "facility_checkins"], RED),
        ("Conversation", ["chat_threads", "chat_messages", "operational_messages"], NAVY),
    ]
    for i, (name, tables, color) in enumerate(clusters):
        x = Inches(0.3 + i * 2.6)
        box(s, x, Inches(1.2), Inches(2.45), Inches(0.5), color, name, 14, True, WHITE, rounded=False)
        box(s, x, Inches(1.7), Inches(2.45), Inches(3.55), SOFT, rounded=False)
        for j, tb in enumerate(tables):
            box(s, x + Inches(0.12), Inches(1.9 + j * 0.8), Inches(2.2), Inches(0.65), WHITE, tb, 13, True, INK, line=color)
    box(
        s,
        Inches(0.3),
        Inches(5.45),
        Inches(12.7),
        Inches(1.45),
        WHITE,
        "UUID PKs  ·  created_at on every table  ·  Proposals reuse appointments (requested + STEP7_PROPOSAL)\n"
        "Conversation reuses chat_threads / chat_messages  ·  Latest ETA is derived, never stored on Shipment\n"
        "Alembic: 888176505d00 + 102c692c1be2 only.  alembic check: no new upgrade operations after freeze.",
        14,
        False,
        INK,
        line=NAVY,
    )

    diagram_slide(
        add, "Diagram  ·  docs/er_diagram.png",
        "Entity-relationship model — 16 frozen tables",
        "er_diagram.png", 0.00, 1.00,
        caption="Shipment and Appointment are the fact hubs. Proposals reuse Appointment. Conversation reuses ChatThread / ChatMessage.",
    )

    # ---------- 12 Steps roadmap ----------
    s = add()
    header(s, "Build", "Steps 1–9 — all implemented, including hardening")
    steps = [
        ("1", "Foundation", "FastAPI, config, health"),
        ("2", "System of record", "16 tables frozen"),
        ("2H", "Hardening", "Indexes + constraints"),
        ("3", "Business APIs", "Read-only, paginated"),
        ("4", "ETA / exceptions", "Facts, not bookings"),
        ("5", "Feasibility", "Pure rule engine"),
        ("6", "Allocation", "Locks + re-check"),
        ("6H", "Hardening", "Allocation safety"),
        ("7", "Proposals", "Show ≠ propose ≠ confirm"),
        ("8", "Conversation", "Allowlisted tools"),
        ("8H", "Hardening", "Adversarial / confirm"),
        ("9", "Scheduling", "Read-only ranking"),
    ]
    for i, (num, title, blurb) in enumerate(steps):
        col, row = i % 6, i // 6
        x = Inches(0.3 + col * 2.15)
        y = Inches(1.25 + row * 2.75)
        color = GOLD if "H" in num or num == "9" else NAVY
        box(s, x, y, Inches(2.05), Inches(2.5), WHITE, "", line=color)
        box(s, x, y, Inches(2.05), Inches(0.7), color, num, 20, True, WHITE, rounded=False)
        textbox(s, x + Inches(0.08), y + Inches(0.85), Inches(1.9), Inches(0.7), title, 14, True, NAVY, PP_ALIGN.CENTER)
        add_lines(s, x + Inches(0.1), y + Inches(1.55), Inches(1.85), Inches(0.8), [blurb], size=12, color=MUTED)

    # ---------- 13 Facts ----------
    s = add()
    header(s, "Step 4", "Record operational facts — delay is not a booking")
    nodes = [
        ("Driver reports\n2 hours late", NAVY),
        ("POST /shipments/{id}\n/eta-updates", BLUE),
        ("Insert eta_updates\nimmutable history", GREEN),
        ("Latest ETA\nderived, not stored\non Shipment", GOLD),
        ("Step 5 reads\nlatest ETA", NAVY),
    ]
    for i, (label, color) in enumerate(nodes):
        box(s, Inches(0.35 + i * 2.6), Inches(1.4), Inches(2.4), Inches(1.55), color, label, 13, True, WHITE)
        if i < len(nodes) - 1:
            arrow_right(s, Inches(2.7 + i * 2.6), Inches(1.98), Inches(0.25), Inches(0.3))
    box(s, Inches(0.35), Inches(3.25), Inches(6.2), Inches(3.55), BLUE_BG, rounded=False)
    textbox(s, Inches(0.55), Inches(3.4), Inches(5.8), Inches(0.4), "ETA UPDATES", 14, True, BLUE)
    add_lines(
        s,
        Inches(0.6),
        Inches(3.95),
        Inches(5.7),
        Inches(2.6),
        [
            "POST /shipments/{id}/eta-updates",
            "GET /shipments/{id}/latest-eta",
            "History is append-only",
            "Conversation tool: record_eta_update",
            "Does not evaluate feasibility",
            "Does not consume a slot",
        ],
        size=15,
        color=INK,
        spacing=6,
    )
    box(s, Inches(6.8), Inches(3.25), Inches(6.15), Inches(3.55), RED_BG, rounded=False)
    textbox(s, Inches(7.0), Inches(3.4), Inches(5.8), Inches(0.4), "DRIVER EXCEPTIONS", 14, True, RED)
    add_lines(
        s,
        Inches(7.05),
        Inches(3.95),
        Inches(5.7),
        Inches(2.6),
        [
            "POST /shipments/{id}/exceptions",
            "open / acknowledged block Step 5 (EXCP-001)",
            "Leave-by 9:30 PM is conversation context,",
            "not automatically a blocking exception",
            "PATCH acknowledges / resolves",
            "Still not a booking",
        ],
        size=15,
        color=INK,
        spacing=6,
    )

    diagram_slide(
        add, "Diagram  ·  docs/eta_exception_flow.png",
        "ETA and exception writes — facts, not bookings",
        "eta_exception_flow.png", 0.00, 0.58,
        caption="Immutable eta_updates history. Latest ETA is derived. Open / acknowledged exceptions block Step 5.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/eta_exception_flow.png",
        "How facts reach feasibility",
        "eta_exception_flow.png", 0.48, 1.00,
        caption="Step 8 tools record facts; Step 5 reads them. Neither write is a capacity commit.",
    )

    # ---------- 14 Feasibility ----------
    s = add()
    header(s, "Step 5", "FeasibilityEngine — pure rules, no database, no LLM")
    add_lines(
        s,
        Inches(0.4),
        Inches(1.15),
        Inches(12.5),
        Inches(0.4),
        ["POST /shipments/{id}/feasibility  →  FeasibilityService  →  FeasibilityEngine.  Outcomes: feasible · not_feasible · not_evaluable"],
        size=14,
        color=MUTED,
    )
    rules = [
        ("SHIP 001–003", "Active, not terminal, destination assigned"),
        ("CARR / DRIV / VEHI", "Actor active; weight/volume when present"),
        ("FACI / APPT", "Facility active; appointment/slot context"),
        ("SLOT 001–004", "Exists, facility match, open, capacity left"),
        ("DOCK 001–005", "Presence, facility, availability, reefer"),
        ("RULE 001–003", "Daily max, hours, dock compatibility"),
        ("ETA 001 / 002", "Latest ETA in window (block); hours warning"),
        ("EXCP-001", "No open or acknowledged exceptions"),
    ]
    for i, (rid, desc) in enumerate(rules):
        col, row = i % 4, i // 4
        x = Inches(0.35 + col * 3.25)
        y = Inches(1.7 + row * 2.15)
        box(s, x, y, Inches(3.1), Inches(1.95), WHITE, "", line=BLUE)
        box(s, x, y, Inches(3.1), Inches(0.5), BLUE, rid, 13, True, WHITE, rounded=False)
        add_lines(s, x + Inches(0.12), y + Inches(0.65), Inches(2.85), Inches(1.15), [desc], size=14, color=INK)
    textbox(
        s,
        Inches(0.4),
        Inches(6.15),
        Inches(12.5),
        Inches(0.75),
        "Capacity-consuming statuses: confirmed and held only.  requested proposals do not consume capacity.\nMissing facts yield not_evaluable — the engine never invents travel times.",
        13,
        False,
        MUTED,
    )

    # ---------- 15 Show propose confirm ----------
    s = add()
    header(s, "Step 7", "Showing is not proposing. Proposing is not confirming.")
    phases = [
        ("1. SHOW", "get_available_options", "Step 5 feasibility\nNumbered options in chat\nNo database write\nCapacity unchanged", BLUE, "“I found these feasible options…”"),
        ("2. PROPOSE", "create_proposal", "Appointment status=requested\nSTEP7_PROPOSAL marker\n30-minute application TTL\nStill not a hold", GOLD, "“Say confirm if you want me to book it.”"),
        ("3. CONFIRM", "accept_proposal", "Revalidate Step 5\nStep 6 row locks\nNew confirmed appointment\nProposal row reconciled", GREEN, "“The appointment is confirmed.”"),
    ]
    for i, (title, tool, body, color, quote) in enumerate(phases):
        x = Inches(0.35 + i * 4.3)
        box(s, x, Inches(1.2), Inches(4.1), Inches(5.7), WHITE, "", line=color)
        box(s, x, Inches(1.2), Inches(4.1), Inches(0.7), color, title, 18, True, WHITE if color != GOLD else NAVY, rounded=False)
        box(s, x + Inches(0.2), Inches(2.05), Inches(3.7), Inches(0.5), SOFT, tool, 13, True, NAVY)
        add_lines(s, x + Inches(0.25), Inches(2.7), Inches(3.6), Inches(2.4), body.split("\n"), size=15, color=INK, spacing=8)
        box(s, x + Inches(0.2), Inches(5.55), Inches(3.7), Inches(1.1), AMBER_BG if i < 2 else GREEN_BG, quote, 12, True, NAVY)
        if i < 2:
            arrow_right(s, Inches(4.35 + i * 4.3), Inches(3.7), Inches(0.28), Inches(0.32))

    diagram_slide(
        add, "Diagram  ·  docs/Show_Propose_Confirm_Sequence.png",
        "Three states as implemented",
        "Show_Propose_Confirm_Sequence.png", 0.00, 0.32,
        caption="Numbered chat options are not a hold. A requested proposal does not consume capacity.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/Show_Propose_Confirm_Sequence.png",
        "Conversation sequence — show / propose / status / confirm",
        "Show_Propose_Confirm_Sequence.png", 0.28, 0.56,
        caption="ASK_STATUS is read-only. Only accept_proposal revalidates and allocates.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/Show_Propose_Confirm_Sequence.png",
        "REST accept internals — same ProposalService",
        "Show_Propose_Confirm_Sequence.png", 0.52, 0.82,
        caption="POST /proposals/{id}/accept: advisory lock → Step 5 → Step 6. World changed → HTTP 409.",
    )

    # ---------- 16 Proposal states ----------
    s = add()
    header(s, "Step 7 lifecycle", "Proposal states — no Proposal table")
    # proposed center
    box(s, Inches(5.15), Inches(2.55), Inches(3.0), Inches(1.15), GOLD, "proposed\n(requested row)", 14, True, NAVY)
    # incoming
    box(s, Inches(5.15), Inches(1.2), Inches(3.0), Inches(0.7), BLUE, "CREATE  (Step 5 feasible)", 12, True, WHITE)
    arrow_down(s, Inches(6.5), Inches(1.92), Inches(0.22), Inches(0.55))
    # terminals
    terms = [
        (Inches(0.4), Inches(2.4), RED, "rejected\nPOST reject"),
        (Inches(0.4), Inches(4.0), SLATE, "expired\nTTL 30 min"),
        (Inches(9.9), Inches(2.4), RGBColor(0x7A, 0x4E, 0x16), "stale\nworld changed"),
        (Inches(9.9), Inches(4.0), GREEN, "confirmed\nseparate row"),
    ]
    for x, y, c, t in terms:
        box(s, x, y, Inches(2.9), Inches(1.2), c, t, 13, True, WHITE)
    # arrows from proposed
    arrow_right(s, Inches(8.2), Inches(2.9), Inches(1.55), Inches(0.28))
    # left arrows - use left arrow shape
    la = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(3.35), Inches(2.9), Inches(1.7), Inches(0.28))
    _fill(la, GOLD)
    arrow_down(s, Inches(6.5), Inches(3.75), Inches(0.22), Inches(0.55))
    box(s, Inches(4.55), Inches(4.35), Inches(4.2), Inches(0.7), SOFT, "accept is in-call only — never persisted as accepted", 12, True, MUTED)
    box(
        s,
        Inches(0.4),
        Inches(5.5),
        Inches(12.5),
        Inches(1.4),
        AMBER_BG,
        "Capacity is consumed only on the confirmed appointment row written by Step 6.\n"
        "Terminal states (rejected / expired / stale / confirmed) have no outbound edges.\n"
        "No PATCH /proposals/{id} status mutation.  No silent retry.  HTTP 409 when the world changed.",
        14,
        False,
        NAVY,
        line=GOLD,
    )

    diagram_slide(
        add, "Diagram  ·  docs/proposal_state_diagram.png",
        "Proposal API states and transitions",
        "proposal_state_diagram.png", 0.00, 0.52,
        caption="accepted is in-call only — never stored. Terminal states have no outbound edges.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/proposal_state_diagram.png",
        "Persistence mapping and accept path",
        "proposal_state_diagram.png", 0.46, 1.00,
        caption="No Proposal table. Confirmed booking is a separate Appointment row written by Step 6.",
    )

    # ---------- 17 Concurrency ----------
    s = add()
    header(s, "Step 6", "Concurrency-safe allocation — one winner for the last unit")
    box(s, Inches(0.35), Inches(1.2), Inches(4.0), Inches(1.0), BLUE, "Request A\nsees remaining capacity = 1", 13, True, WHITE)
    box(s, Inches(4.65), Inches(1.2), Inches(4.0), Inches(1.0), BLUE, "Request B\nsees remaining capacity = 1", 13, True, WHITE)
    box(s, Inches(8.95), Inches(1.2), Inches(3.95), Inches(1.0), SOFT, "Show / propose\ndo not lock", 13, True, NAVY, line=BLUE)
    arrow_down(s, Inches(6.45), Inches(2.25), Inches(0.28), Inches(0.35))
    box(s, Inches(2.4), Inches(2.65), Inches(8.5), Inches(0.7), NAVY, "Step 6 transaction   ·   lock order: shipment  →  slot  →  dock", 15, True, WHITE)
    # lock boxes
    locks = [
        ("1  Advisory lock\nshipment", BLUE),
        ("2  FOR UPDATE\nappointment_slots", GOLD),
        ("3  FOR UPDATE\ndocks", GOLD),
        ("4  Re-run Step 5\ninside the lock", GREEN),
        ("5  Count confirmed+held\nvs slot.capacity", RED),
        ("6  Write confirmed\nsafe_commit", NAVY),
    ]
    for i, (t, c) in enumerate(locks):
        box(s, Inches(0.35 + i * 2.15), Inches(3.55), Inches(2.05), Inches(1.25), c, t, 11, True, WHITE if c != GOLD else NAVY)
        if i < 5:
            arrow_right(s, Inches(2.35 + i * 2.15), Inches(4.0), Inches(0.18), Inches(0.22))
    box(s, Inches(0.35), Inches(5.1), Inches(6.2), Inches(1.7), GREEN_BG, "WINNER\nconfirmed appointment\ncapacity consumed", 16, True, GREEN, line=GREEN)
    box(s, Inches(6.75), Inches(5.1), Inches(6.15), Inches(1.7), RED_BG, "LOSER\nStep 5 / capacity re-check fails\nHTTP 409  ·  stale  ·  no silent retry", 16, True, RED, line=RED)

    diagram_slide(
        add, "Diagram  ·  docs/Concurrency_locking_sequence.png",
        "Lock order — shipment, then slot, then dock",
        "Concurrency_locking_sequence.png", 0.00, 0.28,
        caption="pg_advisory_xact_lock(shipment) then SELECT … FOR UPDATE on appointment_slots and docks.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/Concurrency_locking_sequence.png",
        "Capacity-1 race — two shipments, one slot",
        "Concurrency_locking_sequence.png", 0.26, 0.52,
        caption="Winner is whoever locked the slot row first. Loser re-checks capacity and gets ConflictError / 409.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/Concurrency_locking_sequence.png",
        "Two requested proposals, one confirm",
        "Concurrency_locking_sequence.png", 0.50, 0.78,
        caption="Step 7 accept wraps Step 6. requested rows do not consume capacity until a confirmed row exists.",
    )

    # ---------- 18 Tools ----------
    s = add()
    header(s, "Step 8", "Allowlisted conversation tools — no allocate tool")
    tools = [
        ("get_shipment_status", "READ", BLUE),
        ("record_eta_update", "FACT WRITE", GOLD),
        ("create_driver_exception", "FACT WRITE", GOLD),
        ("evaluate_feasibility", "EVALUATE", GREEN),
        ("get_available_options", "EVALUATE", GREEN),
        ("create_proposal", "PROPOSAL WRITE", RGBColor(0x7A, 0x4E, 0x16)),
        ("get_proposal", "READ", BLUE),
        ("accept_proposal", "CAPACITY COMMIT", RED),
        ("reject_proposal", "PROPOSAL WRITE", RGBColor(0x7A, 0x4E, 0x16)),
        ("request_human_escalation", "ESCALATION FLAG", SLATE),
        ("evaluate_facility_schedule", "EVALUATE (read-only)", GREEN),
    ]
    for i, (name, klass, color) in enumerate(tools):
        col, row = i % 4, i // 4
        if i == 10:
            col, row = 1, 2  # center last? keep sequential - 3 rows of 4, last on row 2 col 2
        x = Inches(0.35 + col * 3.25)
        y = Inches(1.2 + row * 1.85)
        box(s, x, y, Inches(3.1), Inches(1.65), WHITE, "", line=color)
        box(s, x, y, Inches(3.1), Inches(0.45), color, klass, 11, True, WHITE, rounded=False)
        textbox(s, x + Inches(0.08), y + Inches(0.6), Inches(2.95), Inches(0.85), name, 13, True, NAVY, PP_ALIGN.CENTER)
    textbox(s, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.3), "Unknown tool names are rejected.  No eval, exec, or SQL from the model.  Conversation never calls allocate.", 13, True, MUTED)

    # ---------- 19 Hero flow ----------
    s = add()
    header(s, "Hero flow  ·  SH-1024", "Jane Rivera  →  Dallas DC  ·  late inbound, leave-by 9:30 PM")
    turns = [
        ("1", "I'll be two hours late… 8:30 PM", "UPDATE_ETA", "record_eta_update", "Fact write. Not a booking."),
        ("2", "I need to leave by 9:30 PM", "CLARIFICATION", "context only", "leave_by_local in metadata"),
        ("3", "What options do I have?", "ASK_OPTIONS", "get_available_options", "Show. No write."),
        ("4", "The second one works…", "PROPOSE_CHANGE", "create_proposal", "requested. Not confirmed."),
        ("5", "Has it been confirmed?", "ASK_STATUS", "get_proposal", "Read only."),
        ("6", "Confirm it.", "ACCEPT_PROPOSAL", "accept_proposal", "Step 7 → 5 → 6. Booked."),
    ]
    for i, (n, quote, intent, tool, note) in enumerate(turns):
        col, row = i % 3, i // 3
        x = Inches(0.3 + col * 4.3)
        y = Inches(1.2 + row * 2.85)
        box(s, x, y, Inches(4.15), Inches(2.65), WHITE, "", line=NAVY)
        box(s, x, y, Inches(0.55), Inches(2.65), NAVY if i < 5 else GREEN, n, 20, True, WHITE, rounded=False)
        textbox(s, x + Inches(0.7), y + Inches(0.12), Inches(3.3), Inches(0.7), quote, 13, True, INK)
        box(s, x + Inches(0.7), y + Inches(0.9), Inches(3.25), Inches(0.45), SOFT, intent, 11, True, BLUE)
        box(s, x + Inches(0.7), y + Inches(1.45), Inches(3.25), Inches(0.45), AMBER_BG, tool, 11, True, NAVY)
        textbox(s, x + Inches(0.7), y + Inches(2.0), Inches(3.3), Inches(0.5), note, 12, False, MUTED)

    diagram_slide(
        add, "Diagram  ·  docs/driver_conversation_sequence.png",
        "Driver conversation sequence (hero flow)",
        "driver_conversation_sequence.png", 0.00, 0.52,
        caption="Intent → allowlisted tool → existing service. The model never writes SQL or confirms on its own.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/driver_conversation_sequence.png",
        "One-turn internals through the agent",
        "driver_conversation_sequence.png", 0.48, 1.00,
        caption="Guards: missing facts, injection markers, unknown tools, formatter-only wording.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/end_to_end_driver_journey.png",
        "End-to-end journey — SH-1024 at Dallas DC",
        "end_to_end_driver_journey.png", 0.00, 0.55,
        caption="Seeded world first. Delay → leave-by → options → proposal → status → confirm.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/end_to_end_driver_journey.png",
        "From proposal to confirmed appointment",
        "end_to_end_driver_journey.png", 0.48, 1.00,
        caption="scripts/e2e_hero_flow.py against the live API. Stale second confirm cannot silently succeed.",
    )

    # ---------- 20 Step 9 ----------
    s = add()
    header(s, "Step 9  ·  optional", "Read-only facility ranking — never a hold")
    flow = [
        ("facility_id\n+ window", BLUE),
        ("Snapshot\nslots, docks,\nETAs, check-ins", BLUE),
        ("Step 5\neligible pairs", GREEN),
        ("SchedulingEngine\ndeterministic rank", GOLD),
        ("proposed_assignments\n+ unassigned", NAVY),
        ("read_only: true\ncommits_capacity: false", RED),
    ]
    for i, (t, c) in enumerate(flow):
        box(s, Inches(0.3 + i * 2.15), Inches(1.25), Inches(2.05), Inches(1.55), c, t, 12, True, WHITE if c != GOLD else NAVY)
        if i < 5:
            arrow_right(s, Inches(2.3 + i * 2.15), Inches(1.85), Inches(0.18), Inches(0.24))
    ranks = [
        "Confirmed / held appointments are protected and not moved",
        "Earlier facility check-in, then en-route shipments",
        "Lower ETA lateness vs slot end (missing ETA is not fabricated)",
        "Lower early-wait, closer ETA-to-slot-start, then stable IDs",
        "score 0–100 from evaluable ETA metrics; null if ETA missing",
        "Caps: 50 shipments, 100 slots.  No POST /schedule/confirm",
    ]
    for i, t in enumerate(ranks):
        col, row = i % 2, i // 2
        box(s, Inches(0.35 + col * 6.45), Inches(3.1 + row * 1.15), Inches(6.25), Inches(1.0), SOFT, "  " + t, 14, False, INK, PP_ALIGN.LEFT)

    diagram_slide(
        add, "Diagram  ·  docs/Scheduling_architecture.png",
        "Step 9 scheduling architecture",
        "Scheduling_architecture.png", 0.00, 0.55,
        caption="Facility snapshot → Step 5 eligibility → deterministic rank. read_only: true, commits_capacity: false.",
    )
    diagram_slide(
        add, "Diagram  ·  docs/Scheduling_architecture.png",
        "Ranking rules and bounds",
        "Scheduling_architecture.png", 0.48, 1.00,
        caption="Protected confirmed/held rows. No OR-Tools. No POST /schedule/confirm.",
    )

    # ---------- 21 Frontend ----------
    s = add()
    header(s, "Operations console", "React displays API results — it is not a second engine")
    routes = [
        ("/", "Driver Console", "Free-text conversation, intent, tools, status"),
        ("/shipments", "Shipments", "Catalog facts for SH-1024 and peers"),
        ("/appointments", "Appointments", "Proposals and confirmed bookings"),
        ("/facility-schedule", "Facility Schedule", "Step 9 evaluate — labeled read-only"),
        ("/demo", "Demo Scenarios", "Hero messages against seeded world"),
        ("/concurrency", "Concurrency", "Explains the race; does not fire two live confirms"),
    ]
    for i, (path, name, desc) in enumerate(routes):
        col, row = i % 3, i // 3
        x = Inches(0.35 + col * 4.3)
        y = Inches(1.25 + row * 2.5)
        box(s, x, y, Inches(4.1), Inches(2.25), WHITE, "", line=BLUE)
        box(s, x, y, Inches(4.1), Inches(0.55), NAVY, path, 14, True, GOLD, rounded=False)
        textbox(s, x + Inches(0.2), y + Inches(0.7), Inches(3.7), Inches(0.45), name, 18, True, NAVY)
        add_lines(s, x + Inches(0.2), y + Inches(1.2), Inches(3.7), Inches(0.85), [desc], size=14, color=MUTED)
    textbox(s, Inches(0.4), Inches(6.4), Inches(12.5), Inches(0.5), "Vite + React 19 + TypeScript  ·  5 Vitest tests  ·  production build passed  ·  No local feasibility or confirmation math", 13, True, MUTED)

    # ---------- 22 Security ----------
    s = add()
    header(s, "Hardening", "Adversarial and confirmation controls that are actually in the code")
    ctrls = [
        ("Allowlisted tools", "ALLOWED_TOOL_NAMES; unknown → forbidden"),
        ("Strict arguments", "Pydantic extra=forbid; UUID validation"),
        ("Cross-driver bind", "Thread create rejects another driver’s shipment"),
        ("Cross-facility", "Allocation / proposal / Step 9 reject mismatches"),
        ("Prompt injection", "Marker detection; skip irreversible tools"),
        ("Provider privilege", "OpenRouter cannot promote confirm on its own"),
        ("No eval / exec / SQL", "Tools call existing services only"),
        ("Secret stripping", "public_metadata drops keys, prompts, auth"),
        ("Error sanitization", "No traceback or SQLAlchemy leak in chat"),
        ("Stale confirm", "HTTP 409; no silent retry as success"),
        ("Status questions", "ASK_STATUS never calls accept_proposal"),
        ("Classroom auth gap", "No login / JWT — stated limitation"),
    ]
    for i, (t, d) in enumerate(ctrls):
        col, row = i % 4, i // 4
        x = Inches(0.3 + col * 3.25)
        y = Inches(1.2 + row * 1.85)
        box(s, x, y, Inches(3.1), Inches(1.7), WHITE, "", line=NAVY)
        box(s, x, y, Inches(3.1), Inches(0.5), NAVY, t, 12, True, WHITE, rounded=False)
        add_lines(s, x + Inches(0.12), y + Inches(0.65), Inches(2.85), Inches(0.9), [d], size=13, color=INK)

    # ---------- 23 Evidence ----------
    s = add()
    header(s, "Evidence", "What was verified for this submission")
    stats = [
        ("399", "backend tests passed", GREEN),
        ("5", "frontend Vitest tests", BLUE),
        ("16", "frozen domain tables", GOLD),
        ("0", "new tables in Steps 3–9", NAVY),
    ]
    for i, (n, label, c) in enumerate(stats):
        box(s, Inches(0.4 + i * 3.2), Inches(1.25), Inches(3.0), Inches(1.7), c, n, 36, True, WHITE if c != GOLD else NAVY)
        textbox(s, Inches(0.5 + i * 3.2), Inches(2.45), Inches(2.8), Inches(0.4), label, 13, True, WHITE if c != GOLD else NAVY, PP_ALIGN.CENTER)
    box(s, Inches(0.4), Inches(3.2), Inches(6.2), Inches(3.5), SOFT, rounded=False)
    textbox(s, Inches(0.6), Inches(3.35), Inches(5.8), Inches(0.4), "LIVE STACK", 14, True, NAVY)
    add_lines(
        s,
        Inches(0.65),
        Inches(3.9),
        Inches(5.7),
        Inches(2.5),
        [
            "PostgreSQL  ·  localhost:5433",
            "API  ·  127.0.0.1:8010  (not 8000)",
            "UI  ·  127.0.0.1:5173",
            "Health: {status: ok, service: setuhaul}",
            "Hero: scripts/e2e_hero_flow.py",
            "Seed: scripts/seed_ops_demo.py  (SH-1024)",
        ],
        size=15,
        color=INK,
        spacing=6,
    )
    box(s, Inches(6.8), Inches(3.2), Inches(6.1), Inches(3.5), SOFT, rounded=False)
    textbox(s, Inches(7.0), Inches(3.35), Inches(5.7), Inches(0.4), "TEST SPLIT", 14, True, NAVY)
    add_lines(
        s,
        Inches(7.05),
        Inches(3.9),
        Inches(5.7),
        Inches(2.5),
        [
            "API / models: in-memory SQLite",
            "Migration + concurrency: real PostgreSQL",
            "Step 5 / 6 / 7 / 8 hardening suites",
            "Step 6 & 7 concurrency: thread pool, no silent retry",
            "Frontend production build: passed",
            "alembic check: freeze held",
        ],
        size=15,
        color=INK,
        spacing=6,
    )

    # ---------- 24 Demo ----------
    s = add()
    header(s, "Demo", "How to run the walkthrough")
    cmds = [
        ("1", "docker compose up -d", "PostgreSQL on 5433"),
        ("2", "alembic upgrade head", "Frozen Step 2 / 2H schema"),
        ("3", "python scripts/seed_ops_demo.py", "SH-1024, Jane Rivera, Dallas DC"),
        ("4", "uvicorn app.main:app --port 8010", "API on 127.0.0.1:8010"),
        ("5", "npm run dev  (frontend/)", "Console on 127.0.0.1:5173"),
        ("6", "Demo Scenarios or e2e_hero_flow.py", "Delay → options → propose → confirm"),
    ]
    for i, (n, cmd, note) in enumerate(cmds):
        y = Inches(1.2 + i * 0.9)
        box(s, Inches(0.4), y, Inches(0.7), Inches(0.75), GOLD, n, 20, True, NAVY, rounded=False)
        box(s, Inches(1.1), y, Inches(7.6), Inches(0.75), SOFT, "  " + cmd, 15, True, INK, PP_ALIGN.LEFT, rounded=False)
        box(s, Inches(8.8), y, Inches(4.1), Inches(0.75), WHITE, note, 14, False, MUTED, line=RGBColor(0xD0, 0xDA, 0xE4), rounded=False)
    textbox(s, Inches(0.45), Inches(6.65), Inches(12.4), Inches(0.35), "Reset for a clean confirm: docker compose down -v, then steps 1–5.  Re-seed does not undo a confirmed SH-1024.", 13, False, MUTED)

    # ---------- 25 Limitations ----------
    s = add()
    header(s, "Honesty", "Known limitations vs out of scope")
    box(s, Inches(0.35), Inches(1.2), Inches(6.25), Inches(5.7), AMBER_BG, rounded=False)
    textbox(s, Inches(0.55), Inches(1.35), Inches(5.9), Inches(0.4), "LIMITATIONS (designed this way)", 14, True, GOLD)
    lim = [
        "No classroom login / JWT",
        "No email / SMS / push notifications",
        "Escalation is a flag, not a dispatcher SLA",
        "No live travel-time prediction",
        "Proposal TTL is application-side (no expires_at column)",
        "No shipment priority / unload-minutes fields",
        "Step 9 assignment is not a hold",
        "dock_id may be null until proposal / allocation",
        "Conversation does not check the driver in at the gate",
    ]
    add_lines(s, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.7), ["•  " + x for x in lim], size=14, color=INK, spacing=5)
    box(s, Inches(6.8), Inches(1.2), Inches(6.15), Inches(5.7), SOFT, rounded=False)
    textbox(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.4), "OUT OF SCOPE (not this architecture)", 14, True, NAVY)
    oos = [
        "LangChain / LangGraph",
        "Redis / Kafka event bus",
        "OR-Tools / MIP scheduler",
        "National fleet routing",
        "Production identity provider",
        "Notification platform",
        "Human-task inbox / SLA workflow",
        "POST /schedule/confirm",
        "A separate Proposal table",
    ]
    add_lines(s, Inches(7.05), Inches(1.9), Inches(5.7), Inches(4.7), ["•  " + x for x in oos], size=14, color=INK, spacing=5)

    # ---------- 26 Takeaways ----------
    s = add()
    header(s, "Close", "What to remember")
    takes = [
        ("Language vs authority", "The model talks. Steps 5–7 decide and commit."),
        ("Three states", "Show options, write a proposal, then confirm under locks."),
        ("Capacity is scarce", "Only confirmed and held consume a slot. requested does not."),
        ("Revalidate at commit", "A proposal can go stale. The loser gets 409, not a silent book."),
        ("Frozen facts", "16 tables. Later steps reuse Appointment and ChatMessage.metadata."),
        ("Step 9 is optional", "A ranked schedule is advice. Booking still walks 7 → 5 → 6."),
    ]
    for i, (t, d) in enumerate(takes):
        col, row = i % 3, i // 3
        x = Inches(0.35 + col * 4.3)
        y = Inches(1.25 + row * 2.7)
        box(s, x, y, Inches(4.1), Inches(2.45), WHITE, "", line=GOLD)
        box(s, x, y, Inches(4.1), Inches(0.6), NAVY, t, 15, True, WHITE, rounded=False)
        add_lines(s, x + Inches(0.2), y + Inches(0.85), Inches(3.7), Inches(1.4), [d], size=15, color=INK)

    # ---------- 27 Q&A ----------
    s = add()
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _fill(bg, NAVY)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    _fill(accent, GOLD)
    textbox(s, Inches(0.7), Inches(2.1), Inches(12), Inches(0.4), "QUESTIONS", 16, True, GOLD)
    textbox(s, Inches(0.7), Inches(2.55), Inches(12), Inches(1.0), "Thank you", 48, True, WHITE)
    textbox(
        s,
        Inches(0.7),
        Inches(3.7),
        Inches(12),
        Inches(1.2),
        "SETUHAUL  ·  Deterministic appointment coordination\nLLM for language  ·  Python for operational authority",
        18,
        False,
        GOLD_LT,
    )
    box(s, Inches(0.7), Inches(5.3), Inches(3.7), Inches(0.85), RGBColor(0x1A, 0x3A, 0x5C), "Repo: setuhaul\nAPI: 127.0.0.1:8010", 14, True, WHITE)
    box(s, Inches(4.6), Inches(5.3), Inches(3.7), Inches(0.85), RGBColor(0x1A, 0x3A, 0x5C), "Hero: SH-1024\nDallas Distribution Center", 14, True, WHITE)
    box(s, Inches(8.5), Inches(5.3), Inches(4.1), Inches(0.85), RGBColor(0x1A, 0x3A, 0x5C), "Confirm path\n7 → 5 → 6 under locks", 14, True, WHITE)

    total = len(slides)
    for i, sl in enumerate(slides, 1):
        # skip full-bleed title and closing
        if i in (1, total):
            continue
        footer(sl, i, total)

    try:
        prs.save(OUT)
        print(f"Wrote {OUT} ({total} slides)")
    except PermissionError:
        alt = OUT.with_name(OUT.stem + "_updated.pptx")
        prs.save(alt)
        print(f"Original file is open/locked; wrote {alt} ({total} slides)")


if __name__ == "__main__":
    build()
